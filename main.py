from pptx import Presentation
import pyttsx3
import tkinter as tk
from tkinter import ttk
import threading
import time
import os

tts_controller = None

class TTSController:
    def __init__(self, text_to_speak):
        self.text_to_speak = text_to_speak
        self.pause_flag = False
        self.exit_flag = False
        self.speak_thread = None

    def start(self):
        self.speak_thread = threading.Thread(target=self.speak)
        self.speak_thread.start()

    def speak(self):
        bot = pyttsx3.init()
        for text in self.text_to_speak.split('\n'):
            if self.pause_flag:
                while self.pause_flag:
                    time.sleep(1)
            bot.say(text)
            bot.runAndWait()

            if self.exit_flag:
                break

    def pause(self):
        self.pause_flag = True
        print("TTS paused")

    def resume(self):
        self.pause_flag = False
        print("TTS resumed")

    def stop(self):
        self.exit_flag = True
        if self.speak_thread:
            self.speak_thread.join()

def extract_text_from_ppt(ppt_file):
    presentation = Presentation(ppt_file)
    extracted_text = ""
    for slide_number, slide in enumerate(presentation.slides):
        extracted_text += f"\nSlide {slide_number + 1}:\n"
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                extracted_text += f"{shape.text}\n"
    print(extracted_text)
    return extracted_text

def on_pause_button_click():
    tts_controller.pause()

def on_resume_button_click():
    tts_controller.resume()

def on_stop_button_click():
    print("Stop button clicked")
    tts_controller.stop()
    control.destroy()

def on_close():
    tts_controller.stop()
    control.destroy()

def on_back():
    on_close()
    main_gui()

def control_gui():
    global tts_controller
    ppt_file_path = file_path.replace('"','')
    extracted_text_result = extract_text_from_ppt(ppt_file_path)
    
    global control 
    control = tk.Tk()
    control.title("Presentation Reader")

    def startBot():
        tts_controller.start()

    text_area = tk.Text(control, wrap="word", width=50, height=20)
    text_area.insert("1.0", extracted_text_result)
    text_area.config(state="disabled")
    text_area.grid(row=0, column=0, padx=10, pady=10,rowspan=15, columnspan=3)

    start_button = ttk.Button(control, text="Start", command=startBot)
    start_button.grid(row=0, column=3)


    pause_button = ttk.Button(control, text="Pause", command=on_pause_button_click)
    pause_button.grid(row=1, column=3)

    resume_button = ttk.Button(control, text="Resume", command=on_resume_button_click)
    resume_button.grid(row=2, column=3)

    stop_button = ttk.Button(control, text="Exit", command=on_stop_button_click)
    stop_button.grid(row=3, column=3)

    back = ttk.Button(control, text="Back", command=on_back)
    back.grid(row=14, column=3)

    tts_controller = TTSController(extracted_text_result)

    control.protocol("WM_DELETE_WINDOW", on_close)

    control.mainloop()

def main_gui():

    def next_window():
        global file_path
        file_path = file_location_text_area.get("1.0", "end-1c")
        if(os.path.exists(file_path) and ".pptx" in file_path):
            root.destroy()
            control_gui()
        else:
            file_location_invalid_label.grid(row=16, column=2, padx=10, pady=10, sticky="w")


    global root 
    root = tk.Tk()
    root.title("Presentation Reader")
    
    instructions = "Welcome!\nThis program will read a .pptx file and extract the text from it.\nThe text will be converted into an audio which will be played .\n \nHow to use?\n \n1. Place the file in the directory\n2. Then run the program\n3. You will have buttons to pause and resume the audio, the audio will pause only after completing a sentence and when pressed stop the program will exit"
    
    text_area = tk.Text(root, wrap="word", width=50, height=20)
    text_area.insert("1.0", instructions)
    text_area.config(state="disabled")
    text_area.grid(row=0, column=0, padx=10, pady=10, columnspan=3,rowspan=15)

    file_location_ask_label = tk.Label(root, text="Enter the location of the file")
    file_location_ask_label.grid(row=16, column=0, padx=10, pady=10, sticky="w") 

    file_location_invalid_label = tk.Label(root, text="The Pathis is Invalid",fg='red')
     

    file_location_text_area = tk.Text(root, width=40, height=1)
    file_location_text_area.grid(row=17, column=0, sticky="w")
    
    continue_button = tk.Button(root,text="continue",command=next_window)
    continue_button.grid(row=17,column=2,pady=10)
    
    root.mainloop()

main_gui() 