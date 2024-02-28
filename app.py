from pptx import Presentation
import pyttsx3
import tkinter as tk
from tkinter import ttk
import threading
import time
import os
from flask import Flask, render_template, request

app = Flask(__name__, template_folder='templates')
app.config['UPLOAD_FOLDER'] = 'uploads'

tts_controller = None

ALLOWED_EXTENSIONS = {'pptx'}

class TTSController:
    def __init__(self, text_to_speak):
        self.text_to_speak = text_to_speak
        self.pause_flag = False
        self.exit_flag = False
        self.speak_thread = None

    def start(self):
        self.speak_thread = threading.Thread(target=self.speak)
        self.speak_thread.start()

    def speak(self):
        bot = pyttsx3.init()
        for text in self.text_to_speak.split('\n'):
            if self.pause_flag:
                while self.pause_flag:
                    time.sleep(1)
            bot.say(text)
            bot.runAndWait()

            if self.exit_flag:
                break

    def pause(self):
        self.pause_flag = True
        print("TTS paused")

    def resume(self):
        self.pause_flag = False
        print("TTS resumed")

    def stop(self):
        self.exit_flag = True
        if self.speak_thread:
            self.speak_thread.join()

def allowed_file(filename):
    return '.' in filename and filename.rsplit('.', 1)[1].lower() in ALLOWED_EXTENSIONS

def extract_text_from_ppt(ppt_file):
    try:
        presentation = Presentation(ppt_file)
        extracted_text = ""
        for slide_number, slide in enumerate(presentation.slides):
            extracted_text += f"\nSlide {slide_number + 1}:\n"
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    extracted_text += f"{shape.text}\n"
        return extracted_text
    except Exception as e:
        print(f"Error extracting text: {e}")
        return ""

def on_pause_button_click():
    tts_controller.pause()

def on_resume_button_click():
    tts_controller.resume()

def on_stop_button_click():
    print("Stop button clicked")
    tts_controller.stop()

def on_close():
    tts_controller.stop()

def web_control(file_path):
    global tts_controller
    ppt_file_path = file_path.replace('"', '')
    extracted_text_result = extract_text_from_ppt(ppt_file_path)

    tts_controller = TTSController(extracted_text_result)
    threading.Thread(target=tts_controller.start).start()

@app.route('/')
def index():
    return render_template('index.html')

@app.route('/upload', methods=['POST'])
def upload():
    global tts_controller
    uploaded_file = request.files['file']

    if uploaded_file and allowed_file(uploaded_file.filename):
        file_path = os.path.join('uploads', uploaded_file.filename)
        uploaded_file.save(file_path)

        extracted_text_result = extract_text_from_ppt(file_path)
        tts_controller = TTSController(extracted_text_result)
        # Pass extracted_text_result to control.html
        return render_template('control.html', extracted_text=extracted_text_result, file_path=file_path)

    return render_template('index.html', error="Invalid file format or upload failed.")


if __name__ == '__main__':
    app.run(debug=True)

